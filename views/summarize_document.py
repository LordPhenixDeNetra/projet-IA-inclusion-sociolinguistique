import streamlit as st
import os  # Importation pour gérer les chemins de fichiers
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import pandas as pd
from transformers import BartForConditionalGeneration, BartTokenizer

# Initialisation du modèle et du tokenizer pour le résumé
model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')
tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')

# Vérifier si le dossier 'temp' existe et créez-le s'il n'existe pas
if not os.path.exists('temp'):
    os.makedirs('temp')

# Fonction pour extraire le texte de différents formats de fichiers
def extract_text_from_pdf(file_path):
    return extract_text(file_path)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string()

# Fonction pour résumer le texte
def summarize_text(text, max_length=150, min_length=50):
    inputs = tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(inputs, max_length=max_length, min_length=min_length, length_penalty=2.0, num_beams=4, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

# Titre de l'application
st.title("Résumé de Document")

# Upload de fichier
uploaded_file = st.file_uploader("Choisissez un fichier", type=['pdf', 'docx', 'pptx', 'csv'])

if uploaded_file:
    # Sauvegarder temporairement le fichier
    file_path = f"temp/{uploaded_file.name}"
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

    # Extraire le texte selon le type de fichier
    if file_path.endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    elif file_path.endswith('.csv'):
        text = extract_text_from_csv(file_path)
    else:
        st.write("Format de fichier non pris en charge.")
    
    if text:
        # st.write("Texte extrait:")
        # st.write(text)
        
        # Résumer le texte extrait
        summary = summarize_text(text)
        st.write("Résumé:")
        st.write(summary)
    else:
        st.write("Impossible d'extraire du texte du fichier.")
