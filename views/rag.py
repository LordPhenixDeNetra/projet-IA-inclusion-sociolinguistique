import streamlit as st
from transformers import BartForConditionalGeneration, BartTokenizer
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import pandas as pd
import os

# Chargement du modèle et du tokenizer pour la génération de réponse
model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')
tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')

# Fonction de génération de réponse basée sur une requête
def generate_response(question, context):
    input_text = f"Question: {question} Context: {context}"
    inputs = tokenizer(input_text, return_tensors='pt', max_length=1024, truncation=True)
    summary_ids = model.generate(inputs['input_ids'], max_length=150, min_length=50)
    response = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return response

# Fonction de récupération de contexte (simulée)
def retrieve_context(query, documents):
    # Rechercher le document le plus pertinent en fonction de la requête
    for doc in documents:
        if query.lower() in doc.lower():
            return doc
    return "Aucun contexte pertinent trouvé."

# Fonctions d'extraction de texte
def extract_text_from_pdf(file_path):
    return extract_text(file_path)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string()

# Titre de l'application
st.title("Recherche RAG")

# Stockage des documents en mémoire
if "documents" not in st.session_state:
    st.session_state.documents = []

# Upload de fichier
uploaded_file = st.file_uploader("Choisissez un fichier", type=['pdf', 'docx', 'pptx', 'csv'])

if uploaded_file:
    # Sauvegarder temporairement le fichier
    file_path = f"temp/{uploaded_file.name}"
    os.makedirs("temp", exist_ok=True)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

    # Extraire le texte selon le type de fichier
    if file_path.endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    elif file_path.endswith('.csv'):
        text = extract_text_from_csv(file_path)
    else:
        st.write("Format de fichier non pris en charge.")
    
    if text:
        st.session_state.documents.append(text)
        st.write("Texte extrait et stocké avec succès.")
        
# Entrée de l'utilisateur pour la question
query = st.text_input("Entrez votre question ou requête de recherche")

if st.button("Rechercher"):
    if query:
        # Récupérez le contexte pertinent
        context = retrieve_context(query, st.session_state.documents)
        
        # Génération de la réponse
        response = generate_response(query, context)
        
        st.write("Réponse générée:")
        st.write(response)
    else:
        st.write("Veuillez entrer une question ou une requête.")
