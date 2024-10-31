import os
import json
from transformers import BartForConditionalGeneration, BartTokenizer
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import pandas as pd
import streamlit as st

# Chemin du fichier d'index local
INDEX_FILE = "document_index.json"

# Initialisation du modèle BART et du tokenizer
model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')
tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')

# Fonctions pour extraire le texte
def extract_text_from_pdf(file_path):
    return extract_text(file_path)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string()

# Sauvegarder dans un fichier local JSON
def add_file_to_index(file_path):
    if file_path.endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    elif file_path.endswith('.csv'):
        text = extract_text_from_csv(file_path)
    else:
        st.write("Unsupported file format.")
        return

    if text:
        doc = {"file_name": os.path.basename(file_path), "content": text}
        
        if os.path.exists(INDEX_FILE):
            with open(INDEX_FILE, 'r+') as f:
                data = json.load(f)
                data.append(doc)
                f.seek(0)
                json.dump(data, f)
        else:
            with open(INDEX_FILE, 'w') as f:
                json.dump([doc], f)
        
        st.write(f"Document '{os.path.basename(file_path)}' indexed successfully.")

# Fonction de résumé
def summarize_text(text, max_length=150, min_length=50):
    inputs = tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(inputs, max_length=max_length, min_length=min_length, length_penalty=2.0, num_beams=4, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

# Recherche dans les documents indexés
def retrieve_documents(query):
    if not os.path.exists(INDEX_FILE):
        return []
    
    with open(INDEX_FILE, 'r') as f:
        documents = json.load(f)
    
    results = [doc["content"] for doc in documents if query.lower() in doc["content"].lower()]
    return results

# Interface Streamlit
st.title("Document Summarization and Retrieval App")

# Menu de navigation
menu = st.sidebar.selectbox("Choisissez une option", ["Résumé de document", "Recherche RAG", "Résumé de texte"])

# Menu 1 : Résumer un document
if menu == "Résumé de document":
    st.header("Upload Document to Summarize")
    uploaded_file = st.file_uploader("Choose a file", type=['pdf', 'docx', 'pptx', 'csv'])
    
    if uploaded_file is not None:
        file_path = os.path.join("uploads", uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        st.write(f"File '{uploaded_file.name}' uploaded successfully.")
        
        # Résumer le contenu
        if file_path.endswith(('.pdf', '.docx', '.pptx', '.csv')):
            text = add_file_to_index(file_path)
            if text:
                summary = summarize_text(text)
                st.write("Résumé du document :")
                st.write(summary)

# Menu 2 : Recherche RAG
elif menu == "Recherche RAG":
    st.header("Recherche dans les documents indexés")
    query = st.text_input("Enter your search query")
    
    if st.button("Search"):
        if query:
            documents = retrieve_documents(query)
            st.write("Documents pertinents :")
            for doc in documents:
                st.write("- ", doc)
                st.write("\n")

# Menu 3 : Résumé de texte
elif menu == "Résumé de texte":
    st.header("Enter Text to Summarize")
    text_input = st.text_area("Texte à résumer")
    
    if st.button("Summarize"):
        if text_input:
            summary = summarize_text(text_input)
            st.write("Résumé du texte :")
            st.write(summary)
