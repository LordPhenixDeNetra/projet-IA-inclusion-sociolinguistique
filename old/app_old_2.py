import os
import json
from transformers import BartForConditionalGeneration, BartTokenizer
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import pandas as pd
import streamlit as st

# Initialiser le modèle et le tokenizer
model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')
tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')

# Chemin pour sauvegarder les documents localement
INDEX_FILE = "document_index.json"

# Fonction pour extraire du texte à partir de différents types de fichiers
def extract_text_from_pdf(file_path):
    return extract_text(file_path)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string()

# Fonction pour indexer un document dans un fichier JSON
def add_file_to_index(file_path):
    if file_path.endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    elif file_path.endswith('.csv'):
        text = extract_text_from_csv(file_path)
    else:
        st.write("Unsupported file format.")
        return

    if text:
        doc = {"file_name": os.path.basename(file_path), "content": text}
        
        # Sauvegarde dans un fichier JSON
        if os.path.exists(INDEX_FILE):
            with open(INDEX_FILE, 'r+') as f:
                data = json.load(f)
                data.append(doc)
                f.seek(0)
                json.dump(data, f)
        else:
            with open(INDEX_FILE, 'w') as f:
                json.dump([doc], f)
        
        st.write(f"Document '{os.path.basename(file_path)}' indexed successfully.")

# Fonction de résumé
def summarize_text(text, max_length=150, min_length=50):
    inputs = tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(inputs, max_length=max_length, min_length=min_length, length_penalty=2.0, num_beams=4,
                                 early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

# Fonction de recherche locale
def retrieve_documents(query):
    if not os.path.exists(INDEX_FILE):
        return []
    
    with open(INDEX_FILE, 'r') as f:
        documents = json.load(f)
    
    results = [doc["content"] for doc in documents if query.lower() in doc["content"].lower()]
    return results

# Interface Streamlit
def main():
    st.title("Document Summarization and Retrieval")

    # Section de téléchargement de document
    st.header("Upload Document")
    uploaded_file = st.file_uploader("Choose a file", type=['pdf', 'docx', 'pptx', 'csv'])

    if uploaded_file is not None:
        file_path = os.path.join("uploads", uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        st.write(f"File '{uploaded_file.name}' uploaded successfully.")

        # Indexation dans le fichier local
        add_file_to_index(file_path)

    # Section de recherche de document
    st.header("Query Document")
    query = st.text_input("Enter your question or search query")

    if st.button("Search"):
        if query:
            documents = retrieve_documents(query)
            st.write("Relevant Documents:")
            for doc in documents:
                st.write("- ", doc)
                st.write("\n")

    # Section de résumé de texte
    st.header("Summarize Text")
    summarize_text_input = st.text_area("Enter text to summarize")

    if st.button("Summarize"):
        if summarize_text_input:
            summary = summarize_text(summarize_text_input)
            st.write("Summary:")
            st.write(summary)

if __name__ == "__main__":
    main()
