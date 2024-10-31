# main.py
import os
import time
from elasticsearch import Elasticsearch
from transformers import BartForConditionalGeneration, BartTokenizer
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import pandas as pd
import streamlit as st

# main.py
import os
import time


# Initialize Elasticsearch
# es = Elasticsearch([{'host': 'localhost', 'port': 9200, 'scheme': 'http'}])


# Initialisation d'Elasticsearch avec gestion d'erreur
try:
    es = Elasticsearch([{'host': 'localhost', 'port': 9200}])
    if not es.ping():
        raise ValueError("Connexion à Elasticsearch échouée")
    print("Connecté à Elasticsearch avec succès!")
except Exception as e:
    print(f"Erreur de connexion à Elasticsearch: {e}")
    print("Assurez-vous qu'Elasticsearch est installé et démarré sur votre machine")
    print("Vous pouvez le télécharger depuis: https://www.elastic.co/downloads/elasticsearch")
    exit(1)


# Initialize model and tokenizer
model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')
tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')


# Check if Elasticsearch is running and create index if not exist
def check_and_create_index(index_name="documents"):
    if not es.indices.exists(index=index_name):
        es.indices.create(index=index_name)


# Text extraction functions for various file types
def extract_text_from_pdf(file_path):
    return extract_text(file_path)


def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string()


# Document indexing function
def add_file_to_index(file_path, index_name="documents"):
    if file_path.endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    elif file_path.endswith('.csv'):
        text = extract_text_from_csv(file_path)
    else:
        st.write("Unsupported file format.")
        return

    if text:
        doc = {"content": text}
        es.index(index=index_name, body=doc)
        st.write(f"Document '{os.path.basename(file_path)}' indexed successfully.")


# Summarization function
def summarize_text(text, max_length=150, min_length=50):
    inputs = tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(inputs, max_length=max_length, min_length=min_length, length_penalty=2.0, num_beams=4,
                                 early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary


# Document retrieval function
def retrieve_documents(query, index='documents'):
    response = es.search(
        index=index,
        body={"query": {"match": {"content": query}}}
    )
    return [hit['_source']['content'] for hit in response['hits']['hits']]


# Streamlit frontend
def main():
    st.title("Document Summarization and Retrieval")

    # Document Upload Section
    st.header("Upload Document")
    uploaded_file = st.file_uploader("Choose a file", type=['pdf', 'docx', 'pptx', 'csv'])

    if uploaded_file is not None:
        file_path = os.path.join("uploads", uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        st.write(f"File '{uploaded_file.name}' uploaded successfully.")

        # Add to Elasticsearch
        add_file_to_index(file_path)

    # Document Retrieval Section
    st.header("Query Document")
    query = st.text_input("Enter your question or search query")

    if st.button("Search"):
        if query:
            documents = retrieve_documents(query)
            st.write("Relevant Documents:")
            for doc in documents:
                st.write("- ", doc)
                st.write("\n")

    # Text Summarization Section
    st.header("Summarize Text")
    summarize_text_input = st.text_area("Enter text to summarize")

    if st.button("Summarize"):
        if summarize_text_input:
            summary = summarize_text(summarize_text_input)
            st.write("Summary:")
            st.write(summary)


if __name__ == "__main__":
    check_and_create_index()
    main()

